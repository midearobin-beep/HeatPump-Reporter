import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import List, Dict
from datetime import datetime

# Article type code definitions
ARTICLE_TYPE_MAP = {
    "A": "新品发布 / 产品升级",
    "B": "企业战略 / 投资扩产",
    "C": "法规政策 / 补贴",
    "D": "奖项宣传 / 品牌营销",
    "E": "技术趋势 / 行业研究",
    "F": "渠道合作 / 市场进入",
    "G": "财报 / 经营数据",
    "H": "其他",
}

def create_news_ppt(news_items: List[Dict], output_file: str = "Weekly_HeatPump_Report.pptx"):
    prs = Presentation()
    
    # 16:9 Aspect Ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ==========================================
    # COVER SLIDE
    # ==========================================
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.width = prs.slide_width
    title.left = Inches(0)
    title.top = Inches(2.5)
    title.height = Inches(1.5)
    title.text = "全球热泵与 HVAC 行业商业情报日报\n(Commercial Intelligence Daily Briefing)"
    # Force title to 25pt and center alignment
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.size = Pt(25)
        paragraph.alignment = PP_ALIGN.CENTER

    subtitle.width = prs.slide_width
    subtitle.left = Inches(0)
    subtitle.top = Inches(4.2)
    subtitle.height = Inches(1.0)
    subtitle.text = f"报告时间：{datetime.now().strftime('%Y-%m-%d')}"
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.alignment = PP_ALIGN.CENTER

    # ==========================================
    # LEGEND SLIDE: Article Type Definitions
    # ==========================================
    blank_slide_layout = prs.slide_layouts[5]
    legend_slide = prs.slides.add_slide(blank_slide_layout)
    
    legend_title = legend_slide.shapes.title
    legend_title.text = "情报分类索引 (Intelligence Classification)"
    p = legend_title.text_frame.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Left column: Article Type Legend
    txLegend1 = legend_slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5))
    tf_leg1 = txLegend1.text_frame
    tf_leg1.word_wrap = True

    p = tf_leg1.add_paragraph()
    p.text = "文章类型代码"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(50, 50, 50)

    for code, desc in ARTICLE_TYPE_MAP.items():
        p = tf_leg1.add_paragraph()
        p.text = f"  [{code}]  {desc}"
        p.font.size = Pt(14)
        p.space_after = Pt(4)

    # Right column: Rating Legend
    txLegend2 = legend_slide.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.5))
    tf_leg2 = txLegend2.text_frame
    tf_leg2.word_wrap = True

    p = tf_leg2.add_paragraph()
    p.text = "情报等级定义"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(50, 50, 50)

    rating_defs = [
        ("S级  ★★★", "改变格局 / 战略必读", RGBColor(180, 0, 0)),
        ("A级  ★★☆", "高价值 / 重要变动", RGBColor(200, 120, 0)),
        ("B级  ★☆☆", "动态参考 / 行业风向", RGBColor(80, 80, 80)),
    ]
    for label, desc, color in rating_defs:
        p = tf_leg2.add_paragraph()
        p.text = f"  {label}"
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = color
        p = tf_leg2.add_paragraph()
        p.text = f"    {desc}"
        p.font.size = Pt(12)
        p.space_after = Pt(8)

    # Sort news items by continent
    continent_order = ["Europe", "Asia", "North America", "South America", "Oceania", "Africa", "Other"]
    news_items.sort(key=lambda x: continent_order.index(x.get("continent", "Other")) if x.get("continent", "Other") in continent_order else 999)

    # ==========================================
    # NEWS CONTENT SLIDES
    # ==========================================
    for item in news_items:
        # ==========================================
        # SLIDE 1: Executive Summary & Deep Analysis
        # ==========================================
        slide1 = prs.slides.add_slide(blank_slide_layout)
        
        # Title
        title_shape = slide1.shapes.title
        article_type_code = item.get('article_type', 'X')
        article_type_label = ARTICLE_TYPE_MAP.get(article_type_code, article_type_code)
        continent = item.get('continent', 'Other')
        title_text = f"[{continent}] [{article_type_code}·{article_type_label}] {item.get('headline', '无标题')}"
        title_shape.text = title_text
        title_shape.text_frame.word_wrap = False
        p = title_shape.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        # Image (Left side)
        img_bottom = Inches(1.5)
        img_path = item.get('image_path')
        if img_path and os.path.exists(img_path):
            try:
                data_table = item.get("data_table", {})
                pic_height = Inches(2.0) if len(data_table) > 6 else Inches(2.8)
                pic = slide1.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), height=pic_height)
                img_bottom = pic.top + pic.height
            except Exception as e:
                print(f"Error adding image {img_path}: {e}")

        # Text Box (Right side)
        txBox1 = slide1.shapes.add_textbox(Inches(5.6), Inches(1.2), Inches(7.2), Inches(6.0))
        tf1 = txBox1.text_frame
        tf1.word_wrap = True

        # Rating & One-Liner
        rating = item.get('rating', 'N/A')
        rating_color = RGBColor(180, 0, 0) if 'S' in str(rating) else RGBColor(200, 120, 0) if 'A' in str(rating) else RGBColor(80, 80, 80)
        p = tf1.add_paragraph()
        p.text = f"♦ {rating} | {item.get('one_liner', '')}"
        p.font.bold = True
        p.font.color.rgb = rating_color
        p.font.size = Pt(13)
        
        tf1.add_paragraph()
        
        # Summary
        p = tf1.add_paragraph()
        p.text = "【事件摘要】"
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p = tf1.add_paragraph()
        p.text = item.get('summary', '暂无内容。')
        p.font.size = Pt(10)
        p.line_spacing = 1.2

        # Deep Analysis (NEW - the key enrichment)
        deep = item.get('deep_analysis', '')
        if deep:
            tf1.add_paragraph()
            p = tf1.add_paragraph()
            p.text = "【深度分析 & 关键数据】"
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(0, 80, 150)
            p = tf1.add_paragraph()
            p.text = deep
            p.font.size = Pt(10)
            p.line_spacing = 1.3

        # Key Info "Table" (only show fields that have real values)
        key_info = item.get("key_info", {})
        visible_info = {k: v for k, v in key_info.items() if v and str(v).strip() and "未披露" not in str(v)}
        if visible_info:
            tf1.add_paragraph()
            p = tf1.add_paragraph()
            p.text = "【关键参数】"
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(50, 50, 50)
            for k, v in visible_info.items():
                p = tf1.add_paragraph()
                p.text = f" • {k}: {v}"
                p.font.size = Pt(9)

        # Source Links
        tf1.add_paragraph()
        p = tf1.add_paragraph()
        p.text = f"▶ 来源: {item.get('source', '')} | 日期: {item.get('date', '')}"
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(120, 120, 120)
        link = item.get('original_link', '')
        if link:
            p = tf1.add_paragraph()
            p.text = f"▶ 原文链接: {link}"
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor(0, 0, 200)

        # Data Table (Below Left Image)
        data_table = item.get("data_table", {})
        if data_table and isinstance(data_table, dict) and len(data_table) > 0:
            dict_items = [(k, v) for k, v in data_table.items() if v and str(v).strip() and "未披露" not in str(v)]
            if dict_items:
                rows = len(dict_items) + 1
                cols = 2
                table_top = img_bottom + Inches(0.3)
                # Ensure we don't go off the bottom of the slide (7.5 inches limit minus table height buffer)
                if table_top > Inches(4.5):
                    table_top = Inches(4.5)
                
                table_shape = slide1.shapes.add_table(rows, cols, Inches(0.5), table_top, Inches(4.8), Inches(0.25 * rows))
                table = table_shape.table
                
                # Header
                cell0 = table.cell(0, 0)
                cell0.text = "硬核洞察与参数"
                cell0.text_frame.paragraphs[0].font.bold = True
                cell0.text_frame.paragraphs[0].font.size = Pt(10)
                
                cell1 = table.cell(0, 1)
                cell1.text = "数值 / 指标"
                cell1.text_frame.paragraphs[0].font.bold = True
                cell1.text_frame.paragraphs[0].font.size = Pt(10)
                
                # Data Rows
                for i, (k, v) in enumerate(dict_items):
                    c0 = table.cell(i+1, 0)
                    c0.text = str(k)
                    c0.text_frame.paragraphs[0].font.size = Pt(9)
                    
                    c1 = table.cell(i+1, 1)
                    c1.text = str(v)
                    c1.text_frame.paragraphs[0].font.size = Pt(9)

        # ==========================================
        # SLIDE 2: Commercial Intelligence Deep Dive
        # ==========================================
        slide2 = prs.slides.add_slide(blank_slide_layout)
        
        # Subtitle linking to main headline
        title_shape2 = slide2.shapes.title
        title_shape2.text = f"↳ 商业情报洞察: {item.get('headline', '无标题')[:50]}..."
        p2 = title_shape2.text_frame.paragraphs[0]
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.alignment = PP_ALIGN.LEFT
        p2.font.color.rgb = RGBColor(100, 100, 100)

        txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.0), Inches(6.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        # Hidden Signals
        p = tf2.add_paragraph()
        p.text = "【隐藏信号 (Hidden Signals)】"
        p.font.bold = True
        p.font.size = Pt(14)
        for sig in item.get('hidden_signals', []):
            p = tf2.add_paragraph()
            p.text = f"  ⚡ {sig}"
            p.font.size = Pt(12)
        tf2.add_paragraph()

        # Competitor Impact
        p = tf2.add_paragraph()
        p.text = "【竞品影响 (Competitor Impact)】"
        p.font.bold = True
        p.font.size = Pt(14)
        impacts = item.get('competitor_impact', {})
        for region, lines in impacts.items():
            if lines:
                p = tf2.add_paragraph()
                p.text = f"  ♦ {region}:"
                p.font.bold = True
                p.font.size = Pt(11)
                for line in lines:
                    p = tf2.add_paragraph()
                    p.text = f"    → {line}"
                    p.font.size = Pt(11)
        tf2.add_paragraph()

        # Follow-up Suggestions
        p = tf2.add_paragraph()
        p.text = "【建议追踪点 (Follow-ups)】"
        p.font.bold = True
        p.font.size = Pt(14)
        for sug in item.get('suggestions', []):
            p = tf2.add_paragraph()
            p.text = f"  📌 {sug}"
            p.font.size = Pt(12)
        tf2.add_paragraph()

        # Executable Actions
        p = tf2.add_paragraph()
        p.text = "【可执行动作 (Actionable Operations)】"
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(180, 0, 0)
        for act in item.get('actions', []):
            p = tf2.add_paragraph()
            p.text = f"  🔺 {act}"
            p.font.size = Pt(12)
            p.font.bold = True

    prs.save(output_file)
    print(f"PPT 生成完毕，已保存至 {output_file}")

if __name__ == "__main__":
    test_news = [
        {
            "headline": "测试: 大金发布新型高水温热泵",
            "rating": "S级",
            "article_type": "A",
            "one_liner": "大金正试图用不需要管道改造的高温热泵强吃欧洲的传统燃气锅炉存量市场。",
            "summary": "大金在欧洲市场推出第三代针对家庭供暖的高温热泵，最高水温可达70度，无需更换旧型号暖气片，直指英国和德国市场。",
            "deep_analysis": "大金第三代Altherma 3H HT系列最高出水温度提升至70°C（上一代为65°C），COP在35°C水温下达4.56，制热能力覆盖8-16kW。据EHPA数据，2025年欧洲热泵安装量约180万台，同比增长约8%，其中高温机型占比从2024年的12%增至【推测】18%。英国BUS补贴维持£7,500/台，德国BAFA补贴最高可达设备价格的70%。大金此举是针对欧洲约6000万台存量燃气锅炉替换市场的定向打击。",
            "key_info": {
                "国家": "欧洲/英国/德国",
                "品牌": "Daikin",
                "最高水温": "70°C",
                "COP_SCOP": "COP 4.56 @35°C",
                "冷媒": "R32",
                "售价": "文中未披露"
            },
            "hidden_signals": [
                "高温机型抢替代锅炉市场已经成为绝对主流",
                "日系企业开始在欧洲大范围吃补贴红利"
            ],
            "competitor_impact": {
                "对中国制造商": ["警惕其对低端产品的降维打击", "研发需要加快跟进高温冷媒"],
                "对本公司": ["销售资料中加入对此款机型的对比参数"]
            },
            "suggestions": ["关注是否使用R290", "关注最终定价"],
            "actions": ["产品一部：调取大金最新说明书评估参数"],
            "date": "2026-04-10",
            "original_link": "https://example.com"
        }
    ]
    create_news_ppt(test_news, "test_presentation.pptx")
